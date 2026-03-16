#!/usr/bin/env python3
"""
generate_deck.py — HashiCorp × IBM: Vault Agentic IAM Demo Slide Deck

Generates a polished PowerPoint presentation using the HashiCorp FY26 dark theme:
  Background:  #000000 (black)
  Primary:     #FFD714 (Vault yellow)
  Text:        #FFFFFF (white)
  Accent:      #1C1C1C (dark panel) / #2A2A2A (card)
  IBM Blue:    #0F62FE

Usage:
  python3 scripts/generate_deck.py
  # Output: docs/vault-agentic-iam-demo-deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt
import copy
from lxml import etree

# ─────────────────────────────────────────────────────────────────────────────
# Color Palette
# ─────────────────────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x00, 0x00, 0x00)
VAULT_YLW  = RGBColor(0xFF, 0xD7, 0x14)   # Vault yellow — primary brand accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_PANEL = RGBColor(0x1C, 0x1C, 0x1C)   # Slide section backgrounds
CARD_BG    = RGBColor(0x2A, 0x2A, 0x2A)   # Code / callout boxes
MID_GRAY   = RGBColor(0x99, 0x99, 0x99)   # Subtext / de-emphasis
IBM_BLUE   = RGBColor(0x0F, 0x62, 0xFE)   # IBM Carbon blue
IBM_DARK   = RGBColor(0x16, 0x16, 0x16)   # IBM Carbon dark background
GREEN_OK   = RGBColor(0x42, 0xBE, 0x65)   # Success / positive
RED_ALERT  = RGBColor(0xFA, 0x4D, 0x56)   # Danger / problem

# Slide dimensions: Widescreen 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    """Add a blank slide with no placeholder layout."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color=BLACK):
    """Fill the entire slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no border
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    if hasattr(run.font, 'name'):
        run.font.name = "Calibri"
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def set_speaker_notes(slide, notes_text):
    """Set the speaker notes on a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def add_yellow_bar(slide, y=0.55, h=0.06):
    """Add the signature HashiCorp yellow accent bar under titles."""
    add_rect(slide, 0.5, y, 12.3, h, VAULT_YLW)


def add_header(slide, title, subtitle=None, y_title=0.25, y_sub=0.75):
    """Add a consistent slide header: title + optional subtitle."""
    add_text(slide, title,
             x=0.5, y=y_title, w=12.3, h=0.5,
             font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_yellow_bar(slide, y=y_title + 0.45, h=0.045)
    if subtitle:
        add_text(slide, subtitle,
                 x=0.5, y=y_sub, w=12.3, h=0.35,
                 font_size=16, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)


def add_ibm_badge(slide):
    """Add small IBM co-branding text in lower-right corner."""
    add_text(slide, "IBM × HashiCorp",
             x=10.5, y=7.1, w=2.5, h=0.3,
             font_size=9, bold=False, color=MID_GRAY, align=PP_ALIGN.RIGHT)


def add_slide_number(slide, num):
    add_text(slide, str(num),
             x=12.8, y=7.1, w=0.4, h=0.3,
             font_size=9, bold=False, color=MID_GRAY, align=PP_ALIGN.RIGHT)


def add_footer(slide, num):
    add_ibm_badge(slide)
    add_slide_number(slide, num)


def bullet_box(slide, items, x, y, w, h,
               font_size=13, color=WHITE, bullet_color=VAULT_YLW,
               bg_color=None, padding=0.15):
    """Add a bulleted list inside an optional background box."""
    if bg_color:
        add_rect(slide, x, y, w, h, bg_color)

    txBox = slide.shapes.add_textbox(
        Inches(x + padding), Inches(y + padding),
        Inches(w - padding * 2), Inches(h - padding * 2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"▸  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────

def slide_title(prs):
    s = blank_slide(prs)
    fill_bg(s, BLACK)

    # Yellow top bar
    add_rect(s, 0, 0, 13.333, 0.08, VAULT_YLW)

    # IBM blue left accent stripe
    add_rect(s, 0, 0.08, 0.08, 7.42, IBM_BLUE)

    # Main title
    add_text(s, "Vault Agentic IAM",
             x=0.6, y=1.6, w=12.0, h=1.1,
             font_size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    add_text(s, "Eliminating the Confused Deputy Problem",
             x=0.6, y=2.75, w=12.0, h=0.6,
             font_size=26, bold=False, color=VAULT_YLW, align=PP_ALIGN.LEFT)

    # Yellow accent line
    add_rect(s, 0.6, 3.5, 5.0, 0.05, VAULT_YLW)

    # Sub-subtitle
    add_text(s, "How HashiCorp Vault + Microsoft Entra ID deliver\nidentity-aware, ephemeral credentials for AI agents",
             x=0.6, y=3.65, w=9.0, h=0.9,
             font_size=16, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)

    # Co-brand
    add_text(s, "HashiCorp  ×  IBM",
             x=0.6, y=5.0, w=4.0, h=0.5,
             font_size=20, bold=True, color=IBM_BLUE, align=PP_ALIGN.LEFT)

    # URLs block
    add_text(s, "Demo URLs (live):  Web UI → http://localhost:8501  |  Vault UI → http://127.0.0.1:8200/ui",
             x=0.6, y=6.9, w=12.0, h=0.35,
             font_size=9, bold=False, color=MID_GRAY, align=PP_ALIGN.LEFT)

    set_speaker_notes(s, """PRESENTER NOTES — Title Slide

META: Open with the web UI at http://localhost:8501 and the Vault UI at http://127.0.0.1:8200/ui already loaded in separate browser tabs. Have Alice and Bob login windows ready in incognito tabs.

CHALLENGE: Enterprise AI agents increasingly act on behalf of users—but most architectures give agents privileged, long-lived credentials that exceed any individual user's permissions. One compromised agent credential becomes a catastrophic blast radius.

RISK: Without identity-aware credential issuance, agents become confused deputies: they can perform actions no single user is authorized to do, creating audit failures, compliance violations, and security incidents.

SOLUTION: HashiCorp Vault + Microsoft Entra ID enable agents to receive only the credentials the active user is entitled to—short-lived, scoped to the user's group membership, automatically rotated.

OUTCOME: Every agent action is authenticated, authorized to the user's exact permissions, and fully auditable—with zero long-lived credentials in the stack.

TRANSITION: "Let's start by examining the problem this demo solves." → advance to slide 2""")

    add_footer(s, 1)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — THE PROBLEM
# ─────────────────────────────────────────────────────────────────────────────

def slide_problem(prs):
    s = blank_slide(prs)
    fill_bg(s, BLACK)
    add_header(s, "The Problem: The Confused Deputy", "When agents hold the keys to the kingdom")

    # Left: problem diagram box
    add_rect(s, 0.5, 1.35, 5.8, 4.8, DARK_PANEL)
    add_text(s, "Traditional Agent Architecture",
             x=0.7, y=1.5, w=5.4, h=0.4,
             font_size=14, bold=True, color=VAULT_YLW)

    diag = """User Request
      │
      ▼
┌─────────────┐
│    Agent    │ ← holds ADMIN credentials
│  (FastAPI)  │   embedded in config/env
└──────┬──────┘
       │  uses app's own
       │  static credentials
       ▼
┌─────────────┐
│  Database   │ Alice & Bob both
│  (MongoDB)  │ get full access!
└─────────────┘"""
    add_text(s, diag,
             x=0.7, y=1.95, w=5.4, h=3.8,
             font_size=10, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    # Right: consequences
    add_rect(s, 6.7, 1.35, 6.1, 4.8, DARK_PANEL)
    add_text(s, "Consequences",
             x=6.9, y=1.5, w=5.7, h=0.4,
             font_size=14, bold=True, color=RED_ALERT)

    consequences = [
        "Agent acts with MORE privilege than any single user",
        "One compromised API key → full database access",
        "Audit logs show agent identity, not user identity",
        "No way to enforce per-user access controls",
        "Static secrets embedded in containers or env files",
        "Violates least-privilege + separation of duty",
        "Regulatory non-compliance (SOC 2, ISO 27001, HIPAA)",
    ]
    bullet_box(s, consequences, 6.9, 1.95, 5.7, 4.0,
               font_size=12, color=WHITE, bg_color=None)

    add_footer(s, 2)
    set_speaker_notes(s, """PRESENTER NOTES — The Problem

META: Keep this on screen. Do NOT switch to live demo yet.

CHALLENGE: Traditional application architectures give agents service-account credentials with broad permissions. When Alice—a read-only user—asks the agent to list products, the agent uses the same admin credential it uses when Bob—an admin—adds products. The application itself becomes a privileged deputy.

RISK: This creates a critical security gap. If the agent is compromised (prompt injection, supply chain attack, credential leak), attackers inherit full database permissions. There is no per-user boundary. Audit logs are meaningless because every action appears as the agent's identity.

SOLUTION: The solution is to make credential issuance identity-aware—so the agent can only receive credentials that reflect the actual user's permissions at the time of the request.

OUTCOME: Per-user, short-lived credentials eliminate the confused deputy entirely. A compromised agent can only do what the current user is allowed to do—and only for the TTL of those credentials.

TRANSITION: "Here's how Vault solves this." → advance to slide 3""")
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — THE SOLUTION
# ─────────────────────────────────────────────────────────────────────────────

def slide_solution(prs):
    s = blank_slide(prs)
    fill_bg(s, BLACK)
    add_header(s, "The Solution: Vault-Enforced Identity", "Every credential is earned, not assumed")

    # Comparison: Before / After
    add_rect(s, 0.5, 1.35, 5.8, 5.0, DARK_PANEL)
    add_text(s, "❌  Without Vault",
             x=0.7, y=1.45, w=5.4, h=0.4,
             font_size=14, bold=True, color=RED_ALERT)

    before_items = [
        "Static admin credentials in app config",
        "Agent decides authorization logic",
        "All users → same database permissions",
        "Long-lived secrets (months/years)",
        "Secrets in env files or Kubernetes secrets",
        "Manual rotation (or never rotated)",
    ]
    bullet_box(s, before_items, 0.7, 1.9, 5.4, 3.8,
               font_size=12, color=WHITE, bg_color=None)

    add_rect(s, 6.7, 1.35, 6.1, 5.0, DARK_PANEL)
    add_text(s, "✅  With Vault OBO Flow",
             x=6.9, y=1.45, w=5.7, h=0.4,
             font_size=14, bold=True, color=GREEN_OK)

    after_items = [
        "Zero static credentials in the application",
        "Vault enforces authorization via identity groups",
        "Alice → read-only creds  |  Bob → read-write creds",
        "Short-lived credentials (5-min TTL, auto-revoked)",
        "Vault is the only secrets store",
        "Automatic rotation + revocation on expiry",
    ]
    bullet_box(s, after_items, 6.9, 1.9, 5.7, 3.8,
               font_size=12, color=WHITE, bg_color=None)

    # Bottom quote
    add_rect(s, 0.5, 6.5, 12.3, 0.7, CARD_BG)
    add_text(s, "\"The application never decides who gets what access. Vault does—based on the user's verified identity.\"",
             x=0.7, y=6.55, w=12.0, h=0.55,
             font_size=13, bold=False, color=VAULT_YLW,
             align=PP_ALIGN.CENTER, italic=True)

    add_footer(s, 3)
    set_speaker_notes(s, """PRESENTER NOTES — The Solution

META: Keep this on screen. This is the conceptual anchor slide—spend time here before moving to the architecture.

CHALLENGE: The question isn't just "how do we rotate secrets?" It's "how do we make credential issuance identity-aware at runtime, for every request?"

RISK: Partial solutions—like rotating static credentials on a schedule—still give agents more privilege than any user. They reduce secret lifetime but don't solve the authorization problem.

SOLUTION: Vault's JWT auth backend validates the Entra ID OBO token, extracts the user's group membership claims, maps those to Vault policies, and issues database credentials scoped exactly to that user's permissions—for a 5-minute TTL. The application receives only what the user is allowed to receive.

OUTCOME: The application's attack surface collapses to near-zero. Even if an attacker intercepts a credential, it expires in 5 minutes, is scoped to a single user's permissions, and is automatically revoked by Vault.

TRANSITION: "Let me show you the architecture that makes this possible." → advance to slide 4""")
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — ARCHITECTURE OVERVIEW
# ─────────────────────────────────────────────────────────────────────────────

def slide_architecture(prs):
    s = blank_slide(prs)
    fill_bg(s, BLACK)
    add_header(s, "Architecture Overview", "Local stack — zero cloud infrastructure required")

    # Component boxes — row 1
    components = [
        ("products-web\n:8501", "Streamlit UI\nEntra ID OIDC login", IBM_BLUE),
        ("products-agent\n:8001", "FastAPI Agent\nOBO token exchange", VAULT_YLW),
        ("products-mcp\n:8000", "MCP Server\nVault JWT auth", VAULT_YLW),
    ]
    x_positions = [0.5, 4.8, 9.1]
    for (name, desc, color), x in zip(components, x_positions):
        add_rect(s, x, 1.35, 3.7, 1.5, DARK_PANEL)
        add_rect(s, x, 1.35, 3.7, 0.07, color)
        add_text(s, name, x=x+0.15, y=1.5, w=3.4, h=0.6,
                 font_size=14, bold=True, color=color)
        add_text(s, desc, x=x+0.15, y=2.05, w=3.4, h=0.65,
                 font_size=11, bold=False, color=MID_GRAY)

    # Arrows between boxes
    add_text(s, "→", x=4.3, y=1.9, w=0.5, h=0.4,
             font_size=20, bold=True, color=VAULT_YLW, align=PP_ALIGN.CENTER)
    add_text(s, "→", x=8.6, y=1.9, w=0.5, h=0.4,
             font_size=20, bold=True, color=VAULT_YLW, align=PP_ALIGN.CENTER)

    # Row 2: Infrastructure
    infra = [
        ("Vault Enterprise\n:8200", "JWT Auth · Policies\nDynamic DB Secrets", VAULT_YLW),
        ("MongoDB\n:27017", "Products database\ntest.products collection", GREEN_OK),
    ]
    x_positions2 = [2.5, 7.5]
    for (name, desc, color), x in zip(infra, x_positions2):
        add_rect(s, x, 3.3, 3.7, 1.5, DARK_PANEL)
        add_rect(s, x, 3.3, 3.7, 0.07, color)
        add_text(s, name, x=x+0.15, y=3.45, w=3.4, h=0.55,
                 font_size=14, bold=True, color=color)
        add_text(s, desc, x=x+0.15, y=3.95, w=3.4, h=0.65,
                 font_size=11, bold=False, color=MID_GRAY)

    # Vertical arrows
    add_text(s, "↓", x=6.15, y=2.85, w=0.5, h=0.4,
             font_size=20, bold=True, color=VAULT_YLW, align=PP_ALIGN.CENTER)
    add_text(s, "↓", x=6.15, y=2.85, w=0.5, h=0.4,
             font_size=20, bold=True, color=VAULT_YLW, align=PP_ALIGN.CENTER)

    # Vault → MongoDB arrow
    add_text(s, "→  dynamic creds  →", x=6.2, y=3.85, w=1.3, h=0.4,
             font_size=10, bold=False, color=VAULT_YLW, align=PP_ALIGN.CENTER)

    # External services box
    add_rect(s, 0.5, 5.1, 12.3, 1.4, CARD_BG)
    add_text(s, "External Services (required)",
             x=0.7, y=5.15, w=6.0, h=0.35,
             font_size=11, bold=True, color=MID_GRAY)
    external_items = [
        "Microsoft Entra ID — OIDC login + OBO token exchange (identity provider)",
        "AWS Bedrock / Anthropic API — Claude LLM for natural language agent processing",
    ]
    bullet_box(s, external_items, 0.7, 5.45, 12.0, 0.9,
               font_size=11, color=WHITE, bg_color=None)

    add_footer(s, 4)
    set_speaker_notes(s, """PRESENTER NOTES — Architecture Overview

META: Switch to browser tab showing Vault UI: http://127.0.0.1:8200/ui — log in with the root token so the audience can see the namespace, jwt auth mount, and database secrets engine configured by Terraform.

CHALLENGE: Teams need to demonstrate modern AI agent architectures without requiring AWS/cloud infrastructure that's difficult to replicate in customer environments. The full stack must run locally with enterprise-grade security controls.

RISK: Demo environments often use simplified security shortcuts. This demo uses the same patterns you'd use in production: Vault Enterprise namespaces, JWT auth validated against Entra ID JWKS, dynamic credentials with automatic revocation.

SOLUTION: The entire stack runs in nerdctl containers on your laptop. Three pre-built application images from drum0r/ handle the OBO flow, MCP protocol, and Streamlit UI. Vault Enterprise (local) and MongoDB (local) provide the infrastructure. Terraform provisions all Vault configuration automatically.

OUTCOME: Customers see a fully functional, production-equivalent demonstration that they can understand, replicate, and adapt—without any cloud dependencies.

TRANSITION: "Let me walk through the identity flow step by step." → advance to slide 5""")
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — IDENTITY FLOW
# ─────────────────────────────────────────────────────────────────────────────

def slide_identity_flow(prs):
    s = blank_slide(prs)
    fill_bg(s, BLACK)
    add_header(s, "Identity Flow: OBO Token Chain", "How user identity propagates to the database credential")

    steps = [
        ("1", "Browser → products-web",     "User clicks Login → redirected to Entra ID OIDC endpoint",        IBM_BLUE),
        ("2", "products-web → products-agent", "Chat request + user's access token (JWT) forwarded to agent",   IBM_BLUE),
        ("3", "products-agent → Entra ID",  "OBO exchange: jwt-bearer grant → Entra returns OBO token\nrepresenting delegated user identity (groups claim included)", VAULT_YLW),
        ("4", "products-agent → products-mcp", "MCP tool call forwarded with OBO token",                       VAULT_YLW),
        ("5", "products-mcp → Vault",        "JWT login: vault write auth/jwt/login\nVault validates signature via Entra ID JWKS endpoint",            VAULT_YLW),
        ("6", "Vault (internal)",            "Extracts groups claim → matches identity group alias\n→ applies readonly or readwrite policy",             VAULT_YLW),
        ("7", "products-mcp → MongoDB",      "Connects with dynamic credential (5-min TTL)\nVault auto-revokes user on expiry",                         GREEN_OK),
    ]

    row_h = 0.72
    y_start = 1.4
    for i, (num, actor, action, color) in enumerate(steps):
        y = y_start + i * row_h
        # Step number circle
        add_rect(s, 0.5, y, 0.45, 0.5, color)
        add_text(s, num, x=0.5, y=y, w=0.45, h=0.5,
                 font_size=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        # Actor
        add_text(s, actor, x=1.05, y=y, w=3.5, h=0.25,
                 font_size=11, bold=True, color=color)
        # Action
        add_text(s, action, x=1.05, y=y + 0.25, w=11.5, h=0.42,
                 font_size=10, bold=False, color=WHITE)
        # Separator
        if i < len(steps) - 1:
            add_rect(s, 0.5, y + 0.65, 12.3, 0.01, DARK_PANEL)

    add_footer(s, 5)
    set_speaker_notes(s, """PRESENTER NOTES — Identity Flow

META: This is the technical core slide. Keep this visible as you describe each step. You can use this as a reference while the demo runs.

CHALLENGE: Most teams struggle to visualize how a user's identity can propagate through a multi-hop API chain without storing credentials in the application.

RISK: If any step breaks the identity chain—for example, if the agent replaces the OBO token with its own service account token—the entire security model collapses. The application becomes the confused deputy again.

SOLUTION: The OBO (On-Behalf-Of) grant type in OAuth 2.0 is specifically designed for this multi-hop delegation pattern. Entra ID issues a new token at step 3 that represents the user's identity but is scoped for the MCP server audience. Vault then validates this token cryptographically via the JWKS endpoint—never trusting the token blindly.

OUTCOME: By step 7, the database credential MongoDB receives was generated specifically for this user, for this request, and will self-destruct in 5 minutes. No human intervention required. No rotation schedule to manage.

TRANSITION: "Let's run the actual demo. First, let's look at the environment." → advance to slide 6""")
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — DEMO SETUP
# ─────────────────────────────────────────────────────────────────────────────

def slide_demo_setup(prs):
    s = blank_slide(prs)
    fill_bg(s, BLACK)
    add_header(s, "Demo Setup: Environment Ready", "All services running — infrastructure configured by Terraform + Ansible")

    # Left: services status
    add_rect(s, 0.5, 1.35, 5.8, 5.0, DARK_PANEL)
    add_text(s, "Running Services",
             x=0.7, y=1.45, w=5.4, h=0.4,
             font_size=14, bold=True, color=VAULT_YLW)
    services = [
        ("products-web", "http://localhost:8501", GREEN_OK),
        ("products-agent", "http://localhost:8001", GREEN_OK),
        ("products-mcp", "http://localhost:8000", GREEN_OK),
        ("Vault Enterprise", "http://127.0.0.1:8200/ui", GREEN_OK),
        ("MongoDB", "localhost:27017", GREEN_OK),
    ]
    y = 1.95
    for name, url, color in services:
        add_text(s, f"●  {name}", x=0.7, y=y, w=2.8, h=0.35,
                 font_size=12, bold=False, color=color)
        add_text(s, url, x=3.5, y=y, w=2.6, h=0.35,
                 font_size=10, bold=False, color=MID_GRAY)
        y += 0.42

    # Right: demo users
    add_rect(s, 6.7, 1.35, 6.1, 2.4, DARK_PANEL)
    add_text(s, "Demo Users",
             x=6.9, y=1.45, w=5.7, h=0.4,
             font_size=14, bold=True, color=VAULT_YLW)

    # Alice
    add_rect(s, 6.9, 1.9, 5.7, 0.85, CARD_BG)
    add_text(s, "Alice  (Read-Only)", x=7.05, y=1.95, w=5.4, h=0.3,
             font_size=12, bold=True, color=IBM_BLUE)
    add_text(s, "alice@jlevings13gmail699.onmicrosoft.com\ndb-readonly group → readonly Vault policy",
             x=7.05, y=2.25, w=5.4, h=0.45,
             font_size=10, bold=False, color=MID_GRAY)

    # Bob
    add_rect(s, 6.9, 2.85, 5.7, 0.85, CARD_BG)
    add_text(s, "Bob  (Read-Write)", x=7.05, y=2.9, w=5.4, h=0.3,
             font_size=12, bold=True, color=GREEN_OK)
    add_text(s, "bob@jlevings13gmail699.onmicrosoft.com\ndb-readwrite group → readwrite Vault policy",
             x=7.05, y=3.2, w=5.4, h=0.45,
             font_size=10, bold=False, color=MID_GRAY)

    # Vault config summary
    add_rect(s, 6.7, 3.85, 6.1, 2.5, DARK_PANEL)
    add_text(s, "Vault Configuration (via Terraform)",
             x=6.9, y=3.95, w=5.7, h=0.35,
             font_size=13, bold=True, color=VAULT_YLW)
    vault_items = [
        "Namespace: admin/agentic-iam-<suffix>",
        "Auth: JWT (Entra ID JWKS validated)",
        "DB Engine: dynamic MongoDB credentials",
        "TTL: 5 min default · 10 min max",
        "Policies: readonly · readwrite",
    ]
    bullet_box(s, vault_items, 6.9, 4.35, 5.7, 1.9,
               font_size=11, color=WHITE, bg_color=None)

    add_footer(s, 6)
    set_speaker_notes(s, """PRESENTER NOTES — Demo Setup

META: Switch to terminal and run: ./scripts/validate.sh
This confirms all services are healthy. Then switch to browser.

Open these URLs in separate tabs before the demo:
  • Streamlit Web UI:  http://localhost:8501
  • Vault UI:          http://127.0.0.1:8200/ui

Log into Vault UI with the root token (or namespace admin token) so you can show:
  - The 'agentic-iam-xxxx' namespace
  - The JWT auth mount
  - The database secrets engine
  - The readonly and readwrite policies

CHALLENGE: Customers often ask "how do we know the policies are actually enforced?" This slide shows Vault is the enforcement point—not the application.

RISK: Without a live Vault UI walkthrough, the demo can feel abstract. Showing the actual policies and identity groups in Vault makes the security model tangible.

SOLUTION: Terraform provisioned all of this automatically. Zero manual Vault configuration. The reproducibility itself is part of the security story.

OUTCOME: Customers see that Vault's configuration is code-reviewed, version-controlled, and deployable in minutes.

TRANSITION: "Now let's log in as Alice—a read-only user." → advance to slide 7
ACTION: Switch to browser tab showing http://localhost:8501""")
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — ALICE DEMO
# ─────────────────────────────────────────────────────────────────────────────

def slide_alice(prs):
    s = blank_slide(prs)
    fill_bg(s, BLACK)

    # IBM blue accent for Alice
    add_rect(s, 0, 0, 13.333, 0.08, IBM_BLUE)
    add_header(s, "Demo Part 1: Alice — Read-Only User",
               "db-readonly group → readonly Vault policy → read MongoDB credentials")

    # Step-by-step actions
    add_rect(s, 0.5, 1.35, 7.8, 5.0, DARK_PANEL)
    add_text(s, "Live Demo Steps",
             x=0.7, y=1.45, w=7.4, h=0.4,
             font_size=14, bold=True, color=IBM_BLUE)

    steps_alice = [
        ("1. Open App", "Navigate to http://localhost:8501 in browser"),
        ("2. Sign In", "Click Login → Entra ID login page\nEnter: alice@jlevings13gmail699.onmicrosoft.com"),
        ("3. Authenticate", "Enter Alice's password → complete MFA if prompted"),
        ("4. List Products", "Type in chat: List all products\n→ Products appear (read access confirmed)"),
        ("5. Attempt Write", "Type: Add a new product named 'Chair Model X' with price $49.99\n→ Request FAILS — Vault enforced read-only policy"),
        ("6. Observe", "Show the error message: Vault denied write credentials\nThis is the system working correctly"),
    ]
    y = 1.95
    for step, desc in steps_alice:
        add_text(s, step, x=0.7, y=y, w=2.2, h=0.25,
                 font_size=11, bold=True, color=IBM_BLUE)
        add_text(s, desc, x=2.9, y=y, w=5.2, h=0.5,
                 font_size=11, bold=False, color=WHITE)
        y += 0.65

    # Right: expected outcome
    add_rect(s, 8.7, 1.35, 4.1, 5.0, DARK_PANEL)
    add_text(s, "Expected Outcome",
             x=8.9, y=1.45, w=3.7, h=0.4,
             font_size=14, bold=True, color=VAULT_YLW)

    add_rect(s, 8.9, 1.95, 3.7, 0.85, CARD_BG)
    add_text(s, "✅  Can Do", x=9.05, y=2.0, w=3.4, h=0.3,
             font_size=12, bold=True, color=GREEN_OK)
    add_text(s, "Read products (List all products)",
             x=9.05, y=2.3, w=3.4, h=0.45,
             font_size=11, bold=False, color=WHITE)

    add_rect(s, 8.9, 2.95, 3.7, 0.85, CARD_BG)
    add_text(s, "❌  Cannot Do", x=9.05, y=3.0, w=3.4, h=0.3,
             font_size=12, bold=True, color=RED_ALERT)
    add_text(s, "Write / create / delete products",
             x=9.05, y=3.3, w=3.4, h=0.45,
             font_size=11, bold=False, color=WHITE)

    # Vault policy shown
    add_rect(s, 8.9, 3.95, 3.7, 2.0, CARD_BG)
    add_text(s, "Vault Policy: readonly",
             x=9.05, y=4.0, w=3.4, h=0.35,
             font_size=11, bold=True, color=VAULT_YLW)
    policy = """path "database/creds/readonly" {
  capabilities = ["read"]
}"""
    add_text(s, policy, x=9.05, y=4.4, w=3.4, h=1.45,
             font_size=9, bold=False, color=GREEN_OK)

    add_footer(s, 7)
    set_speaker_notes(s, """PRESENTER NOTES — Alice Demo

META: Switch to browser. Use the tab showing http://localhost:8501
If Alice's session is already open, great. If not, open a regular (non-incognito) browser tab.

Alice login credentials:
  Username: alice@jlevings13gmail699.onmicrosoft.com
  Password: (from terraform.tfvars alice_password field)

STEP-BY-STEP ACTIONS:
1. Navigate to http://localhost:8501
2. Click the Login button — you'll be redirected to Microsoft login
3. Enter Alice's email and password
4. In the chat box, type: List all products
   → Wait for the agent to respond with product list
   → If empty first time, run again (MongoDB seed may need a moment)
5. Type: Add a new product named 'Chair Model X' with price $49.99
   → The request will FAIL — this is expected and is the demo working correctly
   → Point out the error message: Vault returned 403 for database/creds/readwrite

CHALLENGE: Customers often ask "but the application could just bypass Vault." This is the moment to show they can't—the application has no database credentials of its own.

RISK: If the write appears to succeed, the Vault readonly policy may not be applied. Abort and check terraform apply output.

SOLUTION: Alice's Entra ID group (db-readonly) maps to the readonly Vault identity group alias, which maps to the readonly policy, which only allows reading database/creds/readonly.

OUTCOME: Alice can read but not write. Vault enforced this—not the application code.

TRANSITION: "Now let's sign in as Bob, who has admin-level permissions." → advance to slide 8
ACTION: Open NEW INCOGNITO window → http://localhost:8501""")
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — BOB DEMO
# ─────────────────────────────────────────────────────────────────────────────

def slide_bob(prs):
    s = blank_slide(prs)
    fill_bg(s, BLACK)

    add_rect(s, 0, 0, 13.333, 0.08, GREEN_OK)
    add_header(s, "Demo Part 2: Bob — Read-Write User",
               "db-readwrite group → readwrite Vault policy → read-write MongoDB credentials")

    add_rect(s, 0.5, 1.35, 7.8, 5.0, DARK_PANEL)
    add_text(s, "Live Demo Steps",
             x=0.7, y=1.45, w=7.4, h=0.4,
             font_size=14, bold=True, color=GREEN_OK)

    steps_bob = [
        ("1. New Incognito Window", "Open NEW incognito/private browser window\nNavigate to http://localhost:8501"),
        ("2. Sign In as Bob", "Click Login → enter: bob@jlevings13gmail699.onmicrosoft.com\nEnter Bob's password → complete MFA if prompted"),
        ("3. List Products", "Type: List all products\n→ Products appear (same as Alice saw)"),
        ("4. Add a Product", "Type: Add a new product named 'Table Model Y' with price $89.99\n→ Operation SUCCEEDS — Vault issued read-write credentials"),
        ("5. Verify Write", "Type: List all products\n→ 'Table Model Y' now appears in the list"),
        ("6. Contrast", "Point out: same application, same agent, different user\n= different Vault credentials = different database permissions"),
    ]
    y = 1.95
    for step, desc in steps_bob:
        add_text(s, step, x=0.7, y=y, w=2.7, h=0.25,
                 font_size=11, bold=True, color=GREEN_OK)
        add_text(s, desc, x=3.4, y=y, w=4.7, h=0.5,
                 font_size=11, bold=False, color=WHITE)
        y += 0.65

    add_rect(s, 8.7, 1.35, 4.1, 5.0, DARK_PANEL)
    add_text(s, "Expected Outcome",
             x=8.9, y=1.45, w=3.7, h=0.4,
             font_size=14, bold=True, color=VAULT_YLW)

    add_rect(s, 8.9, 1.95, 3.7, 0.85, CARD_BG)
    add_text(s, "✅  Can Do", x=9.05, y=2.0, w=3.4, h=0.3,
             font_size=12, bold=True, color=GREEN_OK)
    add_text(s, "Read AND write products\nCreate, update, delete",
             x=9.05, y=2.3, w=3.4, h=0.45,
             font_size=11, bold=False, color=WHITE)

    add_rect(s, 8.9, 2.95, 3.7, 0.85, CARD_BG)
    add_text(s, "Key Point", x=9.05, y=3.0, w=3.4, h=0.3,
             font_size=12, bold=True, color=VAULT_YLW)
    add_text(s, "Same app. Same agent.\nDifferent Vault credentials.",
             x=9.05, y=3.3, w=3.4, h=0.45,
             font_size=11, bold=False, color=WHITE)

    add_rect(s, 8.9, 3.95, 3.7, 2.0, CARD_BG)
    add_text(s, "Vault Policy: readwrite",
             x=9.05, y=4.0, w=3.4, h=0.35,
             font_size=11, bold=True, color=VAULT_YLW)
    policy = """path "database/creds/readwrite" {
  capabilities = ["read"]
}
path "database/creds/readonly" {
  capabilities = ["read"]
}"""
    add_text(s, policy, x=9.05, y=4.4, w=3.4, h=1.45,
             font_size=9, bold=False, color=GREEN_OK)

    add_footer(s, 8)
    set_speaker_notes(s, """PRESENTER NOTES — Bob Demo

META: Use the INCOGNITO window you opened. Do NOT use the same browser session as Alice—the Entra ID session cookie will still be Alice's.

Bob login credentials:
  Username: bob@jlevings13gmail699.onmicrosoft.com
  Password: (from terraform.tfvars bob_password field)

STEP-BY-STEP ACTIONS:
1. In the incognito window, go to http://localhost:8501
2. Click Login → enter Bob's credentials
3. Type in chat: List all products (confirm products visible)
4. Type: Add a new product named 'Table Model Y' with price $89.99
   → This should SUCCEED
5. Type: List all products → confirm Table Model Y appears

KEY TALKING POINT: The application code is identical. The agent is identical. The only difference is which user authenticated and which Entra ID group they belong to. Vault handles all authorization logic externally to the application.

CHALLENGE: Application teams resist externalizing authorization because it feels complex. This demo shows the application literally has zero authorization code—it just calls the MCP server with the user's token.

RISK: If Bob's write fails, check: (1) was readwrite group alias configured in Vault? Run terraform output. (2) Is Bob's Entra ID account in the db-readwrite group?

SOLUTION: The readwrite Vault policy allows reading from database/creds/readwrite, which maps to the MongoDB readWrite role on the test database.

OUTCOME: Bob can read and write. The database credential MongoDB received grants exactly readWrite on the test database—nothing more. Vault issued it for a 5-minute TTL.

TRANSITION: "Let me show you what this looks like inside Vault." → advance to slide 9
ACTION: Switch to Vault UI tab: http://127.0.0.1:8200/ui""")
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — VAULT POLICY MAPPING
# ─────────────────────────────────────────────────────────────────────────────

def slide_policy_mapping(prs):
    s = blank_slide(prs)
    fill_bg(s, BLACK)
    add_header(s, "Inside Vault: Policy Mapping",
               "How Entra ID group membership becomes MongoDB permissions")

    # Flow diagram
    add_rect(s, 0.5, 1.4, 12.3, 4.9, DARK_PANEL)

    # Column headers
    cols = ["Entra ID", "Vault Auth", "Vault Identity", "Vault Policy", "MongoDB Role"]
    x_positions = [0.7, 3.1, 5.5, 8.0, 10.6]
    widths       = [2.2, 2.2, 2.3, 2.3, 2.3]
    for col, x, w in zip(cols, x_positions, widths):
        add_text(s, col, x=x, y=1.5, w=w, h=0.4,
                 font_size=11, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Separator
    add_rect(s, 0.5, 1.93, 12.3, 0.03, CARD_BG)

    # Alice row
    add_text(s, "db-readonly\ngroup", x=0.7, y=2.1, w=2.2, h=0.75,
             font_size=12, bold=True, color=IBM_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, "→", x=2.9, y=2.3, w=0.3, h=0.35,
             font_size=14, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "JWT login\ngroups claim\nextracted", x=3.1, y=2.1, w=2.2, h=0.85,
             font_size=10, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "→", x=5.3, y=2.3, w=0.3, h=0.35,
             font_size=14, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "readonly\nidentity group\nalias match", x=5.5, y=2.1, w=2.3, h=0.85,
             font_size=10, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "→", x=7.8, y=2.3, w=0.3, h=0.35,
             font_size=14, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "readonly\npolicy applied", x=8.0, y=2.1, w=2.3, h=0.75,
             font_size=12, bold=True, color=VAULT_YLW, align=PP_ALIGN.CENTER)
    add_text(s, "→", x=10.4, y=2.3, w=0.25, h=0.35,
             font_size=14, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "read\non test db\n5-min TTL", x=10.6, y=2.1, w=2.3, h=0.85,
             font_size=11, bold=True, color=IBM_BLUE, align=PP_ALIGN.CENTER)

    # Alice label
    add_text(s, "ALICE", x=0.7, y=2.9, w=2.2, h=0.3,
             font_size=9, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)

    add_rect(s, 0.5, 3.2, 12.3, 0.03, CARD_BG)

    # Bob row
    add_text(s, "db-readwrite\ngroup", x=0.7, y=3.4, w=2.2, h=0.75,
             font_size=12, bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)
    add_text(s, "→", x=2.9, y=3.6, w=0.3, h=0.35,
             font_size=14, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "JWT login\ngroups claim\nextracted", x=3.1, y=3.4, w=2.2, h=0.85,
             font_size=10, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "→", x=5.3, y=3.6, w=0.3, h=0.35,
             font_size=14, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "readwrite\nidentity group\nalias match", x=5.5, y=3.4, w=2.3, h=0.85,
             font_size=10, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "→", x=7.8, y=3.6, w=0.3, h=0.35,
             font_size=14, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "readwrite\npolicy applied", x=8.0, y=3.4, w=2.3, h=0.75,
             font_size=12, bold=True, color=VAULT_YLW, align=PP_ALIGN.CENTER)
    add_text(s, "→", x=10.4, y=3.6, w=0.25, h=0.35,
             font_size=14, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, "readWrite\non test db\n5-min TTL", x=10.6, y=3.4, w=2.3, h=0.85,
             font_size=11, bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)

    add_text(s, "BOB", x=0.7, y=4.15, w=2.2, h=0.3,
             font_size=9, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Key insight callout
    add_rect(s, 0.5, 4.6, 12.3, 1.55, CARD_BG)
    add_text(s, "Key Insight: The application code has zero authorization logic. "
             "Vault is the single source of truth for who can do what.",
             x=0.7, y=4.65, w=12.0, h=0.55,
             font_size=14, bold=True, color=VAULT_YLW, align=PP_ALIGN.CENTER)
    add_text(s,
             "Both users hit the same API endpoints. "
             "Vault's identity engine—not application code—determines the credential scope.",
             x=0.7, y=5.2, w=12.0, h=0.7,
             font_size=12, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    add_footer(s, 9)
    set_speaker_notes(s, """PRESENTER NOTES — Policy Mapping

META: Switch to Vault UI: http://127.0.0.1:8200/ui
Navigate to: Namespace (agentic-iam-xxxx) → Access → Identity → Groups
Show the 'readonly' and 'readwrite' external identity groups.
Navigate to: Policies → Show the readonly and readwrite HCL policies.
Navigate to: Secrets → database → Roles — show the readonly and readwrite roles with 5-min TTL.

CHALLENGE: Customers want to know how Vault "knows" which user is which. The JWKS validation is the critical piece.

RISK: Without seeing this in Vault UI, customers may assume the application is doing the group lookup. The Vault UI makes it undeniable that this is Vault-enforced.

SOLUTION: When the MCP server sends vault write auth/jwt/login role=default jwt=<OBO_TOKEN>, Vault fetches the JWKS endpoint from Entra ID (https://login.microsoftonline.com/<tenant>/discovery/v2.0/keys), validates the token signature, extracts the groups claim (which contains Entra ID group Object IDs), and matches them against identity group aliases.

OUTCOME: This mapping is version-controlled in Terraform. Changes to group assignments are infrastructure changes—reviewable, auditable, reproducible.

TRANSITION: "Let me now show the credential lifecycle—how short-lived credentials work." → advance to slide 10""")
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — SHORT-LIVED CREDENTIALS
# ─────────────────────────────────────────────────────────────────────────────

def slide_credentials(prs):
    s = blank_slide(prs)
    fill_bg(s, BLACK)
    add_header(s, "Short-Lived Credentials: The Last Line of Defense",
               "5-minute TTL · Auto-revoked · No rotation schedule needed")

    # Timeline visual
    add_rect(s, 0.5, 1.5, 12.3, 0.08, VAULT_YLW)  # timeline bar

    events = [
        (0.5, "t=0\nRequest", "User sends\nchat message", VAULT_YLW),
        (3.2, "t~1s\nOBO Exchange", "Entra ID issues\nOBO token", IBM_BLUE),
        (5.9, "t~2s\nVault Login", "JWT validated\nCreds issued", VAULT_YLW),
        (8.5, "t~3s\nDB Operation", "Query executes\nwith dynamic creds", GREEN_OK),
        (11.0, "t=5min\nAuto-Revoke", "Vault revokes\nMongoDB user", RED_ALERT),
    ]
    for x, label, desc, color in events:
        add_rect(s, x - 0.08, 1.44, 0.16, 0.22, color)
        add_text(s, label, x=x - 0.7, y=1.7, w=1.4, h=0.55,
                 font_size=9, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, desc, x=x - 0.7, y=2.3, w=1.4, h=0.55,
                 font_size=9, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Three benefits columns
    cols_data = [
        ("Zero Blast Radius", VAULT_YLW, [
            "Credentials expire in 5 minutes",
            "Vault auto-revokes MongoDB user",
            "No rotation schedule to manage",
            "Intercepted creds are useless in seconds",
        ]),
        ("Full Auditability", IBM_BLUE, [
            "Every credential issuance logged in Vault",
            "Tied to specific user + timestamp",
            "Vault audit log is tamper-evident",
            "Compliance-ready out of the box",
        ]),
        ("Zero App Secrets", GREEN_OK, [
            "No database passwords in containers",
            "No .env files with real credentials",
            "No Kubernetes secrets with DB creds",
            "Application starts with no secrets at all",
        ]),
    ]
    x_start = 0.5
    col_w = 4.0
    for i, (title, color, items) in enumerate(cols_data):
        x = x_start + i * (col_w + 0.1)
        add_rect(s, x, 3.0, col_w, 3.8, DARK_PANEL)
        add_rect(s, x, 3.0, col_w, 0.07, color)
        add_text(s, title, x=x+0.15, y=3.1, w=col_w - 0.3, h=0.45,
                 font_size=14, bold=True, color=color)
        bullet_box(s, items, x+0.15, 3.6, col_w - 0.3, 3.1,
                   font_size=12, color=WHITE, bg_color=None)

    add_footer(s, 10)
    set_speaker_notes(s, """PRESENTER NOTES — Short-Lived Credentials

META: Stay on this slide — no live demo action needed here. This is the "so what" slide that converts technical features into business value.

CHALLENGE: Security teams understand the value, but business stakeholders ask: "How does this reduce risk in practice? What does a 5-minute TTL actually prevent?"

RISK: Without concrete impact framing, short TTLs sound like an implementation detail. The answer: it transforms a stolen credential from a persistent threat into an expiring liability.

SOLUTION: Walk through the "what if" scenario: Suppose a prompt injection attack causes the agent to exfiltrate credentials. With static credentials, the attacker has permanent database access until someone notices and manually rotates. With Vault dynamic credentials, the attacker has at most 5 minutes of access—and Vault's audit log will show exactly when and which credentials were issued.

OUTCOME: This is the foundation for Zero Trust data access. Every database operation is authenticated, authorized, and time-bounded. No long-lived secrets anywhere in the stack.

REAL-WORLD IMPACT: Organizations using Vault dynamic secrets have:
  - Reduced MTTR for credential-related incidents by 80%+
  - Eliminated credential rotation toil entirely
  - Achieved SOC 2 Type II compliance faster (audit-ready logs)

TRANSITION: "Let me now summarize the enterprise value of this pattern." → advance to slide 11""")
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — BUSINESS VALUE
# ─────────────────────────────────────────────────────────────────────────────

def slide_business_value(prs):
    s = blank_slide(prs)
    fill_bg(s, BLACK)
    add_header(s, "Enterprise Value: Why This Pattern Matters",
               "Security, compliance, and developer velocity — not trade-offs")

    quadrants = [
        ("Security", VAULT_YLW, 0.5, 1.4, [
            "Eliminates confused deputy attack surface",
            "Zero persistent credentials in applications",
            "Short-lived creds limit breach impact window",
            "Vault audit log: every access, every credential",
            "JWKS-validated identity—no token forgery possible",
        ]),
        ("Compliance", IBM_BLUE, 6.95, 1.4, [
            "Least-privilege enforced at infrastructure level",
            "SOC 2, ISO 27001, HIPAA audit-ready logs",
            "User identity in every database operation",
            "Policy-as-code: version-controlled, reviewable",
            "Automated credential lifecycle = no compliance gaps",
        ]),
        ("Developer Experience", GREEN_OK, 0.5, 4.3, [
            "Zero secrets management in application code",
            "Same code deploys across all environments",
            "No secret rotation outages or on-call pages",
            "Access control changes via Terraform PR—not code",
            "New users provisioned via Entra ID group—instant access",
        ]),
        ("Operational", MID_GRAY, 6.95, 4.3, [
            "Full-stack deploys in minutes via Terraform + Ansible",
            "Reproducible: 'terraform destroy && terraform apply'",
            "Vault namespaces isolate team environments",
            "Works locally and in production—same config",
            "No cloud dependencies for core security pattern",
        ]),
    ]
    for title, color, x, y, items in quadrants:
        add_rect(s, x, y, 5.9, 2.75, DARK_PANEL)
        add_rect(s, x, y, 5.9, 0.07, color)
        add_text(s, title, x=x+0.15, y=y+0.1, w=5.6, h=0.4,
                 font_size=14, bold=True, color=color)
        bullet_box(s, items, x+0.15, y+0.55, 5.6, 2.0,
                   font_size=11, color=WHITE, bg_color=None)

    add_footer(s, 11)
    set_speaker_notes(s, """PRESENTER NOTES — Business Value

META: No screen switch needed. This is a conversation slide. Pause here and invite questions.

CHALLENGE: CISOs and CIOs must balance security investment against developer productivity. Traditional PAM and secrets rotation tools are expensive to operate and create friction for developers.

RISK: Organizations that defer secrets management often accumulate technical debt (hard-coded credentials in 47 microservices) that creates a security crisis during audits or after a breach.

SOLUTION: Vault + Entra ID creates a single control plane for all identity and access decisions. Developers write application code that has zero secrets awareness. Security teams manage policy centrally. Compliance teams get automatic audit logs.

OUTCOME: HashiCorp customers report:
  - 60-80% reduction in secrets-related security incidents
  - Elimination of manual credential rotation (saving hundreds of engineering hours/year)
  - Audit preparation time reduced from weeks to hours
  - IBM Security's Zero Trust framework alignment: identity as the new perimeter

TALKING POINTS FOR IBM CO-SELL:
  - HashiCorp Vault integrates with IBM Security Verify (ISAM) as an additional identity provider
  - IBM Cloud Pak for Security can ingest Vault audit logs for SIEM correlation
  - IBM Consulting practices can deploy this pattern for enterprise customers

TRANSITION: "Here's how to get started with this pattern in your organization." → advance to slide 12""")
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — CALL TO ACTION
# ─────────────────────────────────────────────────────────────────────────────

def slide_cta(prs):
    s = blank_slide(prs)
    fill_bg(s, BLACK)

    # Yellow top bar
    add_rect(s, 0, 0, 13.333, 0.08, VAULT_YLW)
    add_rect(s, 0, 0.08, 0.08, 7.42, IBM_BLUE)

    add_text(s, "Get Started",
             x=0.6, y=0.5, w=12.0, h=0.8,
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_rect(s, 0.6, 1.35, 5.0, 0.05, VAULT_YLW)

    # Next steps
    add_rect(s, 0.5, 1.55, 6.1, 4.5, DARK_PANEL)
    add_text(s, "Next Steps",
             x=0.7, y=1.65, w=5.7, h=0.4,
             font_size=16, bold=True, color=VAULT_YLW)
    next_steps = [
        "Run this demo in your own environment (GitHub link below)",
        "Schedule a Vault Architecture Workshop with IBM Consulting",
        "Assess your current secrets sprawl with HashiCorp Secrets Audit",
        "Pilot dynamic secrets in one application team",
        "Expand to HCP Vault Dedicated for production",
    ]
    bullet_box(s, next_steps, 0.7, 2.1, 5.7, 3.7,
               font_size=12, color=WHITE, bg_color=None)

    # Resources box
    add_rect(s, 6.9, 1.55, 5.9, 4.5, DARK_PANEL)
    add_text(s, "Field Resources",
             x=7.1, y=1.65, w=5.5, h=0.4,
             font_size=16, bold=True, color=VAULT_YLW)
    resources = [
        "Demo Source:  github.com/Jlevings/demo-vault-agentic-iam",
        "Vault Docs:   developer.hashicorp.com/vault",
        "HCP Vault:    portal.cloud.hashicorp.com",
        "Vault OBO Blog: hashicorp.com/blog (search 'confused deputy')",
        "IBM × HashiCorp: ibm.com/hashicorp",
    ]
    bullet_box(s, resources, 7.1, 2.1, 5.5, 3.7,
               font_size=11, color=WHITE, bg_color=None)

    # Demo live links
    add_rect(s, 0.5, 6.2, 12.3, 0.9, CARD_BG)
    add_text(s,
             "Demo Live:  Web UI → http://localhost:8501  |  "
             "Vault UI → http://127.0.0.1:8200/ui  |  "
             "Alice: alice@jlevings13gmail699.onmicrosoft.com  |  Bob: bob@jlevings13gmail699.onmicrosoft.com",
             x=0.7, y=6.3, w=12.0, h=0.65,
             font_size=10, bold=False, color=VAULT_YLW, align=PP_ALIGN.CENTER)

    add_footer(s, 12)
    set_speaker_notes(s, """PRESENTER NOTES — Call to Action

META: End the live demo here. Close the browser tabs showing Alice and Bob sessions. Return to this slide.

CHALLENGE: Customers leave demos excited but unsure where to start. The path from demo to production feels long.

RISK: Without a clear next step, the demo becomes a one-time event with no follow-through. The goal is a committed next action—a workshop, a pilot, a proof-of-concept.

SOLUTION: HashiCorp and IBM offer a structured path:
  1. Secrets Audit → identify where static credentials exist today
  2. Vault Architecture Workshop → design the target state
  3. Pilot → implement dynamic secrets for one team/app
  4. Scale → HCP Vault Dedicated or Vault Enterprise cluster

OUTCOME: Customers who complete a Vault pilot typically expand to 3-5 use cases within 6 months: dynamic secrets, PKI as a service, SSH cert auth, KV secrets, OIDC federation.

IBM CO-SELL: This demo is a lead-in for IBM Security Services' Zero Trust practice and HashiCorp's Channel+ partner program.

CLOSING STATEMENT:
"What you saw today is running live on a laptop. The same architecture scales to HCP Vault with thousands of applications. The security model doesn't change—it just gets bigger. Let's talk about where you want to start."

META: Offer to email the GitHub link and schedule a follow-up workshop.
GitHub: https://github.com/Jlevings/demo-vault-agentic-iam""")
    return s


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_identity_flow(prs)
    slide_demo_setup(prs)
    slide_alice(prs)
    slide_bob(prs)
    slide_policy_mapping(prs)
    slide_credentials(prs)
    slide_business_value(prs)
    slide_cta(prs)

    out_path = "docs/vault-agentic-iam-demo-deck.pptx"
    prs.save(out_path)
    print(f"✅  Deck saved to {out_path}")
    print(f"    Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
